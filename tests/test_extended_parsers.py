"""
Unit tests for extended document extractors in services/attachments.py.
Tests dynamic file generation (.docx, .xlsx, .pptx, .csv, .md, .py, .js, .json, .html, .css)
and error handling for corrupt files.
"""

import io
import pytest
import docx
import openpyxl
import pptx

from services.attachments import (
    _extract_text_from_docx,
    _extract_text_from_xlsx,
    _extract_text_from_pptx,
    _extract_text_from_csv,
    _extract_text_from_markdown,
    _extract_text_from_json,
    _extract_text_from_html,
    _extract_text_from_code,
    extract_attachment_content,
    is_supported_attachment,
)


@pytest.mark.asyncio
async def test_extract_docx_success_and_corrupt():
    """Test dynamic DOCX file creation, paragraph & table extraction, and corrupt handling."""
    doc = docx.Document()
    doc.add_heading("Document Title", level=1)
    doc.add_paragraph("First paragraph content for RAG testing.")

    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Row1 Cell1"
    table.cell(1, 1).text = "Row1 Cell2"

    bio = io.BytesIO()
    doc.save(bio)
    docx_bytes = bio.getvalue()

    # Direct function test
    text = _extract_text_from_docx(docx_bytes)
    assert text is not None
    assert "Document Title" in text
    assert "First paragraph content for RAG testing." in text
    assert "Header A | Header B" in text
    assert "Row1 Cell1 | Row1 Cell2" in text

    # extract_attachment_content test
    details = await extract_attachment_content(docx_bytes, "test_doc.docx")
    assert details is not None
    assert details["file_type"] == "docx"
    assert details["file_ext"] == ".docx"
    assert details["page_or_sheet_count"] >= 1
    assert "Document Title" in details["text"]

    # Corrupt file test
    corrupt_text = _extract_text_from_docx(b"CORRUPTED_DOCX_BYTES")
    assert corrupt_text is None

    corrupt_details = await extract_attachment_content(b"CORRUPTED_DOCX_BYTES", "corrupt.docx")
    assert corrupt_details is None


@pytest.mark.asyncio
async def test_extract_xlsx_success_and_corrupt():
    """Test dynamic XLSX workbook creation, sheet & row extraction, and corrupt handling."""
    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Sales"
    ws1.append(["Product", "Revenue"])
    ws1.append(["Widget A", 1000])
    ws1.append(["Widget B", 2500])

    ws2 = wb.create_sheet(title="Expenses")
    ws2.append(["Category", "Amount"])
    ws2.append(["Rent", 500])

    bio = io.BytesIO()
    wb.save(bio)
    xlsx_bytes = bio.getvalue()

    # Direct function test
    text = _extract_text_from_xlsx(xlsx_bytes)
    assert text is not None
    assert "--- Sheet: Sales ---" in text
    assert "Product | Revenue" in text
    assert "Widget A | 1000" in text
    assert "--- Sheet: Expenses ---" in text
    assert "Rent | 500" in text

    # extract_attachment_content test
    details = await extract_attachment_content(xlsx_bytes, "financials.xlsx")
    assert details is not None
    assert details["file_type"] == "xlsx"
    assert details["file_ext"] == ".xlsx"
    assert details["page_or_sheet_count"] == 2
    assert "Widget B | 2500" in details["text"]

    # Corrupt file test
    corrupt_text = _extract_text_from_xlsx(b"INVALID_XLSX_DATA")
    assert corrupt_text is None

    corrupt_details = await extract_attachment_content(b"INVALID_XLSX_DATA", "corrupt.xlsx")
    assert corrupt_details is None


@pytest.mark.asyncio
async def test_extract_pptx_success_and_corrupt():
    """Test dynamic PPTX presentation creation, slide & shape extraction, and corrupt handling."""
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[0]  # Title slide layout

    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Presentation Overview"
    slide1.placeholders[1].text = "Subtitle: Key RAG Features"

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Architecture Diagram"

    bio = io.BytesIO()
    prs.save(bio)
    pptx_bytes = bio.getvalue()

    # Direct function test
    text = _extract_text_from_pptx(pptx_bytes)
    assert text is not None
    assert "--- Slide 1: Presentation Overview ---" in text
    assert "Subtitle: Key RAG Features" in text
    assert "--- Slide 2: Architecture Diagram ---" in text

    # extract_attachment_content test
    details = await extract_attachment_content(pptx_bytes, "slides.pptx")
    assert details is not None
    assert details["file_type"] == "pptx"
    assert details["file_ext"] == ".pptx"
    assert details["page_or_sheet_count"] == 2
    assert "Presentation Overview" in details["text"]

    # Corrupt file test
    corrupt_text = _extract_text_from_pptx(b"NOT_A_VALID_PPTX")
    assert corrupt_text is None

    corrupt_details = await extract_attachment_content(b"NOT_A_VALID_PPTX", "bad.pptx")
    assert corrupt_details is None


@pytest.mark.asyncio
async def test_extract_csv_success():
    """Test CSV file extraction and formatting."""
    csv_data = b"ID,Name,Role\n101,Alice,Engineer\n102,Bob,Manager\n"

    text = _extract_text_from_csv(csv_data)
    assert text is not None
    assert "ID | Name | Role" in text
    assert "101 | Alice | Engineer" in text
    assert "102 | Bob | Manager" in text

    details = await extract_attachment_content(csv_data, "data.csv")
    assert details is not None
    assert details["file_type"] == "csv"
    assert details["file_ext"] == ".csv"
    assert details["page_or_sheet_count"] == 1


@pytest.mark.asyncio
async def test_extract_markdown_success():
    """Test Markdown file extraction preserving headers and list formatting."""
    md_data = (
        b"# Main Title\n\n"
        b"Introduction paragraph.\n\n"
        b"## Features\n\n"
        b"- Feature 1\n"
        b"- Feature 2\n"
    )

    text = _extract_text_from_markdown(md_data, "README.md")
    assert text is not None
    assert "# Main Title" in text
    assert "## Features" in text
    assert "- Feature 1" in text

    details = await extract_attachment_content(md_data, "doc.md")
    assert details is not None
    assert details["file_type"] == "markdown"
    assert details["file_ext"] == ".md"


@pytest.mark.asyncio
async def test_extract_code_files():
    """Test Python, JS, and CSS code file extraction wrapped in syntax blocks."""
    # Python
    py_data = b"def calculate_total(a, b):\n    return a + b\n"
    py_details = await extract_attachment_content(py_data, "script.py")
    assert py_details is not None
    assert py_details["file_type"] == "code"
    assert "```python" in py_details["text"]
    assert "def calculate_total" in py_details["text"]

    # JavaScript
    js_data = b"const fetchData = async () => {\n  console.log('fetching');\n};\n"
    js_details = await extract_attachment_content(js_data, "app.js")
    assert js_details is not None
    assert js_details["file_type"] == "code"
    assert "```javascript" in js_details["text"]
    assert "console.log('fetching')" in js_details["text"]

    # CSS
    css_data = b".header { background-color: #333; color: white; }\n"
    css_details = await extract_attachment_content(css_data, "styles.css")
    assert css_details is not None
    assert css_details["file_type"] == "code"
    assert "```css" in css_details["text"]
    assert ".header {" in css_details["text"]


@pytest.mark.asyncio
async def test_extract_json_success_and_corrupt():
    """Test JSON formatting and invalid JSON handling."""
    json_data = b'{"project": "RagDiscord", "version": 2, "active": true}'

    text = _extract_text_from_json(json_data)
    assert text is not None
    assert "```json" in text
    assert '"project": "RagDiscord"' in text
    assert '"version": 2' in text

    details = await extract_attachment_content(json_data, "config.json")
    assert details is not None
    assert details["file_type"] == "json"

    # Corrupt JSON test
    corrupt_json = _extract_text_from_json(b'{"incomplete_json": ')
    assert corrupt_json is None

    corrupt_details = await extract_attachment_content(b'{"incomplete_json": ', "bad.json")
    assert corrupt_details is None


@pytest.mark.asyncio
async def test_extract_html_success():
    """Test HTML extraction using BeautifulSoup, filtering script/style tags."""
    html_data = (
        b"<!DOCTYPE html><html>"
        b"<head><title>Test Page</title><style>body { margin: 0; }</style></head>"
        b"<body>"
        b"<h1>HTML Header</h1>"
        b"<script>alert('xss');</script>"
        b"<p>Clean text inside HTML body.</p>"
        b"</body></html>"
    )

    text = _extract_text_from_html(html_data)
    assert text is not None
    assert "HTML Header" in text
    assert "Clean text inside HTML body." in text
    assert "alert('xss')" not in text
    assert "margin: 0" not in text

    details = await extract_attachment_content(html_data, "index.html")
    assert details is not None
    assert details["file_type"] == "html"
    assert "HTML Header" in details["text"]


def test_is_supported_attachment():
    """Test filename extension verification for all supported formats."""
    assert is_supported_attachment("document.docx") is True
    assert is_supported_attachment("data.xlsx") is True
    assert is_supported_attachment("deck.pptx") is True
    assert is_supported_attachment("table.csv") is True
    assert is_supported_attachment("readme.md") is True
    assert is_supported_attachment("main.py") is True
    assert is_supported_attachment("script.js") is True
    assert is_supported_attachment("config.json") is True
    assert is_supported_attachment("index.html") is True
    assert is_supported_attachment("styles.css") is True
    assert is_supported_attachment("manual.pdf") is True
    assert is_supported_attachment("unknown.exe") is False
