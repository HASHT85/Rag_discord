"""
Traitement des pièces jointes Discord.

Télécharge et extrait le texte des fichiers joints aux messages
Discord (PDF, DOCX, XLSX, PPTX, CSV, MD, Code, JSON, HTML, images) pour l'indexation RAG.
"""

import io
import json
import logging
from pathlib import Path
from typing import Any

import aiohttp
import discord
import fitz  # PyMuPDF

logger = logging.getLogger(__name__)

# ─────────────────────────────────────────────
#  Constantes
# ─────────────────────────────────────────────

# Taille maximale de fichier autorisée (10 Mo)
_MAX_FILE_SIZE: int = 10 * 1024 * 1024

# Extensions de fichiers texte supportées
_TEXT_EXTENSIONS: set[str] = {
    ".txt", ".md", ".py", ".json", ".csv", ".log",
    ".xml", ".yaml", ".yml", ".html", ".htm", ".css", ".js", ".ts",
}

# Extensions de documents riches supportées
_DOC_EXTENSIONS: set[str] = {
    ".docx", ".xlsx", ".pptx",
}

# Extensions d'images supportées (décrites par le LLM vision)
_IMAGE_EXTENSIONS: set[str] = {
    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp",
}

# Toutes les extensions supportées (texte + docs + PDF + images)
_SUPPORTED_EXTENSIONS: set[str] = _TEXT_EXTENSIONS | _DOC_EXTENSIONS | {".pdf"} | _IMAGE_EXTENSIONS


def is_supported_attachment(filename: str) -> bool:
    """
    Vérifie si l'extension du fichier est supportée pour l'extraction de texte.

    Args:
        filename: Nom du fichier (avec extension).

    Returns:
        True si le fichier est supporté, False sinon.
    """
    ext = Path(filename).suffix.lower()
    supported = ext in _SUPPORTED_EXTENSIONS
    logger.debug(
        "📎 Fichier '%s' (ext=%s) — supporté : %s",
        filename, ext, supported,
    )
    return supported


async def _download_attachment(attachment: discord.Attachment) -> bytes | None:
    """
    Télécharge le contenu d'une pièce jointe Discord de manière asynchrone.

    Args:
        attachment: Objet Attachment de discord.py.

    Returns:
        Contenu du fichier en bytes, ou None en cas d'erreur ou de dépassement de taille.
    """
    if attachment.size and attachment.size > _MAX_FILE_SIZE:
        logger.warning(
            "⚠️ Fichier '%s' trop volumineux (%d octets, max=%d).",
            attachment.filename, attachment.size, _MAX_FILE_SIZE,
        )
        return None

    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(attachment.url) as response:
                if response.status != 200:
                    logger.error(
                        "❌ Échec du téléchargement de '%s' (HTTP %d).",
                        attachment.filename, response.status,
                    )
                    return None

                data = await response.read()

                if len(data) > _MAX_FILE_SIZE:
                    logger.warning(
                        "⚠️ Fichier '%s' trop volumineux après téléchargement (%d octets).",
                        attachment.filename, len(data),
                    )
                    return None

                logger.debug(
                    "📥 Fichier '%s' téléchargé (%d octets).",
                    attachment.filename, len(data),
                )
                return data

    except aiohttp.ClientError as e:
        logger.error(
            "❌ Erreur réseau lors du téléchargement de '%s' : %s",
            attachment.filename, e,
        )
        return None


def _decode_bytes(data: bytes) -> str | None:
    """Décode les octets bruts en UTF-8 puis latin-1."""
    for encoding in ("utf-8", "latin-1"):
        try:
            return data.decode(encoding)
        except (UnicodeDecodeError, ValueError):
            continue
    return None


def _extract_text_from_pdf(data: bytes) -> str | None:
    """Extrait le texte d'un fichier PDF."""
    try:
        doc = fitz.open(stream=data, filetype="pdf")
        pages_text: list[str] = []

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text("text")
            if text and text.strip():
                pages_text.append(text.strip())

        doc.close()

        if not pages_text:
            logger.warning("⚠️ Aucun texte extrait du PDF.")
            return None

        full_text = "\n\n".join(pages_text)
        logger.info(
            "📄 Texte PDF extrait : %d page(s), %d caractères.",
            len(pages_text), len(full_text),
        )
        return full_text

    except Exception as e:
        logger.error("❌ Erreur lors de l'extraction PDF : %s", e)
        return None


def _extract_text_from_docx(data: bytes) -> str | None:
    """Extrait le texte et les tableaux d'un fichier Word (.docx)."""
    try:
        import docx
        doc = docx.Document(io.BytesIO(data))
        parts: list[str] = []

        for paragraph in doc.paragraphs:
            p_text = paragraph.text.strip()
            if p_text:
                parts.append(p_text)

        for table in doc.tables:
            for row in table.rows:
                cells_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells_text:
                    parts.append(" | ".join(cells_text))

        if not parts:
            logger.warning("⚠️ Aucun texte extrait du fichier DOCX.")
            return None

        full_text = "\n\n".join(parts)
        logger.info("📄 Texte DOCX extrait : %d élément(s), %d caractères.", len(parts), len(full_text))
        return full_text

    except Exception as e:
        logger.error("❌ Erreur lors de l'extraction DOCX : %s", e)
        return None


def _extract_text_from_xlsx(data: bytes) -> str | None:
    """Extrait les feuilles, en-têtes et lignes d'un fichier Excel (.xlsx)."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
        sheet_sections: list[str] = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows_text: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                non_empty = [str(cell).strip() for cell in row if cell is not None and str(cell).strip() != ""]
                if non_empty:
                    rows_text.append(" | ".join(non_empty))

            if rows_text:
                sheet_content = f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows_text)
                sheet_sections.append(sheet_content)

        if not sheet_sections:
            logger.warning("⚠️ Aucun texte extrait du fichier XLSX.")
            return None

        full_text = "\n\n".join(sheet_sections)
        logger.info("📊 Texte XLSX extrait : %d feuille(s), %d caractères.", len(sheet_sections), len(full_text))
        return full_text

    except Exception as e:
        logger.error("❌ Erreur lors de l'extraction XLSX : %s", e)
        return None


def _extract_text_from_pptx(data: bytes) -> str | None:
    """Extrait les diapositives, titres et contenus d'un fichier PowerPoint (.pptx)."""
    try:
        import pptx
        prs = pptx.Presentation(io.BytesIO(data))
        slide_sections: list[str] = []

        for i, slide in enumerate(prs.slides, start=1):
            slide_lines: list[str] = []
            title_text = ""
            if slide.shapes.title and slide.shapes.title.text:
                title_text = slide.shapes.title.text.strip()

            header = f"--- Slide {i}: {title_text} ---" if title_text else f"--- Slide {i} ---"

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text and text != title_text:
                            slide_lines.append(text)

            content_str = "\n".join(slide_lines)
            if title_text or slide_lines:
                section = f"{header}\n{content_str}".strip() if content_str else header
                slide_sections.append(section)

        if not slide_sections:
            logger.warning("⚠️ Aucun texte extrait du fichier PPTX.")
            return None

        full_text = "\n\n".join(slide_sections)
        logger.info("📙 Texte PPTX extrait : %d diapositive(s), %d caractères.", len(slide_sections), len(full_text))
        return full_text

    except Exception as e:
        logger.error("❌ Erreur lors de l'extraction PPTX : %s", e)
        return None


def _extract_text_from_csv(data: bytes) -> str | None:
    """Extrait et formate les tableaux de données CSV."""
    try:
        import csv
        text_content = _decode_bytes(data)
        if not text_content:
            return None

        reader = csv.reader(io.StringIO(text_content))
        formatted_rows: list[str] = []
        for row in reader:
            non_empty = [cell.strip() for cell in row if cell.strip()]
            if non_empty:
                formatted_rows.append(" | ".join(cell.strip() for cell in row))

        if not formatted_rows:
            logger.warning("⚠️ Aucun texte extrait du fichier CSV.")
            return None

        full_text = "\n".join(formatted_rows)
        logger.info("📈 Texte CSV extrait : %d ligne(s), %d caractères.", len(formatted_rows), len(full_text))
        return full_text

    except Exception as e:
        logger.error("❌ Erreur lors de l'extraction CSV : %s", e)
        return None


def _extract_text_from_markdown(data: bytes, filename: str = "") -> str | None:
    """Extrait le texte Markdown en préservant en-têtes et structures de listes."""
    return _extract_text_from_text_file(data, filename)


def _extract_text_from_code(data: bytes, filename: str, ext: str) -> str | None:
    """Enrobe le code dans des blocs de syntaxe ```lang ... ```."""
    raw_text = _extract_text_from_text_file(data, filename)
    if not raw_text:
        return None

    lang_map = {
        ".py": "python",
        ".js": "javascript",
        ".ts": "typescript",
        ".css": "css",
        ".html": "html",
        ".json": "json",
        ".sh": "bash",
        ".sql": "sql",
        ".cpp": "cpp",
        ".c": "c",
    }
    lang = lang_map.get(ext.lower(), "")
    return f"```{lang}\n{raw_text}\n```"


def _extract_text_from_json(data: bytes) -> str | None:
    """Extrait et formate du JSON en chaîne prettifiée."""
    try:
        text_content = _decode_bytes(data)
        if not text_content:
            return None

        parsed_json = json.loads(text_content)
        formatted_json = json.dumps(parsed_json, indent=2, ensure_ascii=False)
        return f"```json\n{formatted_json}\n```"

    except Exception as e:
        logger.error("❌ Erreur lors de l'extraction JSON : %s", e)
        return None


def _extract_text_from_html(data: bytes) -> str | None:
    """Extrait le texte lisible d'un document HTML via BeautifulSoup."""
    try:
        from bs4 import BeautifulSoup
        text_content = _decode_bytes(data)
        if not text_content:
            return None

        soup = BeautifulSoup(text_content, "html.parser")
        for s in soup(["script", "style", "head", "title", "meta"]):
            s.decompose()

        body_or_html = soup.body if soup.body else soup
        extracted_text = body_or_html.get_text(separator="\n", strip=True)
        if not extracted_text or not extracted_text.strip():
            logger.warning("⚠️ Aucun texte extrait du fichier HTML.")
            return None

        return extracted_text.strip()

    except Exception as e:
        logger.error("❌ Erreur lors de l'extraction HTML : %s", e)
        return None


def _extract_text_from_text_file(data: bytes, filename: str) -> str | None:
    """Décode le contenu d'un fichier texte brut (UTF-8 puis latin-1)."""
    text = _decode_bytes(data)
    if text and text.strip():
        logger.debug(
            "📝 Fichier '%s' décodé (%d car.).",
            filename, len(text),
        )
        return text.strip()

    logger.warning("⚠️ Impossible de décoder le fichier '%s'.", filename)
    return None


async def extract_attachment_content(data: bytes, filename: str) -> dict[str, Any] | None:
    """
    Extrait le texte et les métadonnées d'un contenu de pièce jointe selon son extension.

    Args:
        data: Contenu binaire du fichier.
        filename: Nom du fichier.

    Returns:
        Dictionnaire avec 'text', 'file_type', 'file_ext', 'page_or_sheet_count',
        ou None si non supporté ou si l'extraction échoue.
    """
    ext = Path(filename).suffix.lower()
    if not is_supported_attachment(filename):
        logger.info("ℹ️ Format non supporté pour '%s' (ext=%s).", filename, ext)
        return None

    text: str | None = None
    file_type: str = "text"
    page_or_sheet_count: int = 1

    if ext == ".pdf":
        file_type = "pdf"
        text = _extract_text_from_pdf(data)
        if text:
            try:
                doc = fitz.open(stream=data, filetype="pdf")
                page_or_sheet_count = len(doc)
                doc.close()
            except Exception:
                page_or_sheet_count = 1

    elif ext == ".docx":
        file_type = "docx"
        text = _extract_text_from_docx(data)
        if text:
            try:
                import docx
                doc = docx.Document(io.BytesIO(data))
                page_or_sheet_count = max(1, len(doc.paragraphs))
            except Exception:
                page_or_sheet_count = 1

    elif ext == ".xlsx":
        file_type = "xlsx"
        text = _extract_text_from_xlsx(data)
        if text:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
                page_or_sheet_count = len(wb.sheetnames)
            except Exception:
                page_or_sheet_count = 1

    elif ext == ".pptx":
        file_type = "pptx"
        text = _extract_text_from_pptx(data)
        if text:
            try:
                import pptx
                prs = pptx.Presentation(io.BytesIO(data))
                page_or_sheet_count = len(prs.slides)
            except Exception:
                page_or_sheet_count = 1

    elif ext == ".csv":
        file_type = "csv"
        text = _extract_text_from_csv(data)
        page_or_sheet_count = 1

    elif ext == ".md":
        file_type = "markdown"
        text = _extract_text_from_markdown(data, filename)
        page_or_sheet_count = 1

    elif ext == ".json":
        file_type = "json"
        text = _extract_text_from_json(data)
        page_or_sheet_count = 1

    elif ext in (".html", ".htm"):
        file_type = "html"
        text = _extract_text_from_html(data)
        page_or_sheet_count = 1

    elif ext in (".py", ".js", ".ts", ".css"):
        file_type = "code"
        text = _extract_text_from_code(data, filename, ext)
        page_or_sheet_count = 1

    elif ext in _IMAGE_EXTENSIONS:
        file_type = "image"
        text = await _describe_image_with_llm(data, filename)
        page_or_sheet_count = 1

    elif ext in _TEXT_EXTENSIONS:
        file_type = "text"
        text = _extract_text_from_text_file(data, filename)
        page_or_sheet_count = 1

    else:
        logger.warning("⚠️ Extension '%s' non gérée.", ext)
        return None

    if not text or not text.strip():
        logger.warning("⚠️ Aucun texte extrait de '%s'.", filename)
        return None

    return {
        "text": text,
        "file_type": file_type,
        "file_ext": ext,
        "page_or_sheet_count": page_or_sheet_count,
    }


async def extract_text_from_attachment(
    attachment: discord.Attachment,
) -> str | None:
    """
    Télécharge et extrait le texte d'une pièce jointe Discord.
    """
    data = await _download_attachment(attachment)
    if data is None:
        return None

    content_info = await extract_attachment_content(data, attachment.filename)
    if content_info:
        logger.info(
            "✅ Texte extrait de '%s' : %d caractères.", attachment.filename, len(content_info["text"])
        )
        return content_info["text"]

    logger.warning("⚠️ Aucun texte extrait de '%s'.", attachment.filename)
    return None


async def extract_attachment_details(
    attachment: discord.Attachment,
) -> dict[str, Any] | None:
    """
    Télécharge et extrait le texte et les métadonnées d'une pièce jointe Discord.
    """
    data = await _download_attachment(attachment)
    if data is None:
        return None

    return await extract_attachment_content(data, attachment.filename)


async def _describe_image_with_llm(data: bytes, filename: str) -> str | None:
    """Décrit une image via le LLM vision (Gemini Flash)."""
    from services.openrouter_client import describe_image

    try:
        description = await describe_image(data, filename)
        if description:
            logger.info(
                "🖼️ Image '%s' décrite par le LLM (%d caractères).",
                filename, len(description),
            )
            return f"[Image : {filename}]\n{description}"
        return None
    except Exception as e:
        logger.error("❌ Erreur description image '%s' : %s", filename, e)
        return None
